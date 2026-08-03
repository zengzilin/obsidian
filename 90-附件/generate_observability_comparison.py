from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import landscape
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from reportlab.platypus import Paragraph, Table, TableStyle
from reportlab.pdfgen import canvas


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "Dynatrace_vs_ELK_Prometheus_Comparison.pptx"
PDF_PATH = OUT_DIR / "Dynatrace_vs_ELK_Prometheus_Comparison.pdf"
PREVIEW_PATH = OUT_DIR / "Dynatrace_vs_ELK_Prometheus_Comparison_preview.png"

TITLE = "Observability Platform Comparison"
SUBTITLE = "Company context: ELK for APM and logs; Prometheus for infrastructure and business monitoring"

ROWS = [
    ("Log Search & Ad Hoc Analysis", "Strong; mature full-text search and flexible troubleshooting workflows", "Strong; query workflows and ingestion costs require validation"),
    ("APM & Distributed Tracing", "Good; broad service coverage, but current trace volume creates substantial ES pressure", "Stronger automatic instrumentation, discovery, and contextual analysis"),
    ("Service Topology", "Depends on agent coverage and consistent service metadata", "More mature automatic topology and dependency discovery"),
    ("Root Cause Analysis", "Manual correlation across Kibana, Prometheus, and alerts", "Stronger automated correlation, anomaly detection, and causal analysis"),
    ("Kubernetes Monitoring", "Prometheus is cloud-native and already established in the environment", "Highly automated, with infrastructure-to-application context"),
    ("Business Metrics", "Mature PromQL queries, exporters, dashboards, and alerting rules", "Requires integration or migration of existing metrics and rules"),
    ("Advanced Log Analytics", "Major Elastic strength; familiar search and investigation experience", "Capable, but DQL adoption and retention economics must be assessed"),
    ("Platform Operations", "High: Kafka, Logstash, Elasticsearch, lifecycle, shard, and capacity management", "Lower infrastructure burden; more responsibility moves to the vendor"),
    ("Cost Model", "Infrastructure, storage, and engineering labor; trace volume is the main cost driver", "Commercial licensing and data consumption; requires a measured PoC"),
    ("Data Ownership & Control", "High, with self-managed storage and retention policies", "Lower; depends on the selected SaaS or managed deployment model"),
    ("Vendor Lock-in", "Medium", "High"),
    ("Migration Risk", "None for the current operating model", "Very high for a full replacement"),
]

NAVY = RGBColor(22, 34, 51)
TEAL = RGBColor(0, 128, 125)
CORAL = RGBColor(226, 92, 76)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(39, 48, 58)
MUTED = RGBColor(99, 111, 123)
PALE = RGBColor(244, 247, 249)
LINE = RGBColor(210, 218, 224)


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_cell(cell, text, size, color, bold=False, fill=None):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.025)
    cell.margin_bottom = Inches(0.02)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_text(slide, 0.48, 0.23, 8.4, 0.42, TITLE, 24, NAVY, True)
    add_text(slide, 0.49, 0.67, 11.95, 0.30, SUBTITLE, 9.5, MUTED)

    tag = slide.shapes.add_shape(5, Inches(11.45), Inches(0.25), Inches(1.38), Inches(0.38))
    tag.fill.solid()
    tag.fill.fore_color.rgb = CORAL
    tag.line.fill.background()
    add_text(slide, 11.45, 0.25, 1.38, 0.38, "DECISION VIEW", 8, WHITE, True, PP_ALIGN.CENTER)

    table_shape = slide.shapes.add_table(len(ROWS) + 1, 3, Inches(0.48), Inches(1.08), Inches(12.35), Inches(5.92))
    table = table_shape.table
    table.columns[0].width = Inches(2.15)
    table.columns[1].width = Inches(5.05)
    table.columns[2].width = Inches(5.15)
    table.rows[0].height = Inches(0.45)
    for row_index in range(1, len(ROWS) + 1):
        table.rows[row_index].height = Inches(0.456)

    headers = ["DIMENSION", "CURRENT ELK + PROMETHEUS", "FULL DYNATRACE ADOPTION"]
    for idx, header in enumerate(headers):
        set_cell(table.cell(0, idx), header, 9.2, WHITE, True, NAVY if idx == 0 else TEAL)

    for row_idx, values in enumerate(ROWS, start=1):
        base_fill = WHITE if row_idx % 2 else PALE
        for col_idx, value in enumerate(values):
            fill = base_fill
            text_color = INK
            bold = col_idx == 0
            if values[0] == "Migration Risk" and col_idx == 2:
                fill = RGBColor(253, 236, 233)
                text_color = CORAL
                bold = True
            elif values[0] == "Vendor Lock-in" and col_idx == 2:
                text_color = CORAL
                bold = True
            set_cell(table.cell(row_idx, col_idx), value, 7.8 if col_idx else 8.1, text_color, bold, fill)

    for row in table.rows:
        for cell in row.cells:
            cell.border_left = None
            cell.border_right = None

    add_text(slide, 0.49, 7.12, 12.3, 0.20, "Assessment focus: capability fit, operating burden, data economics, and migration risk.", 8, MUTED)
    prs.save(PPTX_PATH)


def get_pdf_font():
    candidates = [
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/calibri.ttf"),
    ]
    for font_path in candidates:
        if font_path.exists():
            pdfmetrics.registerFont(TTFont("DeckFont", str(font_path)))
            return "DeckFont"
    return "Helvetica"


def build_pdf():
    page_size = landscape((13.333333 * inch, 7.5 * inch))
    pdf = canvas.Canvas(str(PDF_PATH), pagesize=page_size)
    width, height = page_size
    font_name = get_pdf_font()

    pdf.setFillColor(colors.HexColor("#FFFFFF"))
    pdf.rect(0, 0, width, height, fill=1, stroke=0)
    pdf.setFillColor(colors.HexColor("#00807D"))
    pdf.rect(0, 0, 0.16 * inch, height, fill=1, stroke=0)

    pdf.setFillColor(colors.HexColor("#162233"))
    pdf.setFont(font_name, 24)
    pdf.drawString(0.48 * inch, height - 0.58 * inch, TITLE)
    pdf.setFillColor(colors.HexColor("#636F7B"))
    pdf.setFont(font_name, 9.5)
    pdf.drawString(0.49 * inch, height - 0.91 * inch, SUBTITLE)

    pdf.setFillColor(colors.HexColor("#E25C4C"))
    pdf.roundRect(11.45 * inch, height - 0.63 * inch, 1.38 * inch, 0.38 * inch, 0.12 * inch, fill=1, stroke=0)
    pdf.setFillColor(colors.white)
    pdf.setFont(font_name, 8)
    pdf.drawCentredString(12.14 * inch, height - 0.49 * inch, "DECISION VIEW")

    styles = {
        "header": ParagraphStyle("header", fontName=font_name, fontSize=9.2, leading=10.2, textColor=colors.white, alignment=TA_LEFT),
        "dim": ParagraphStyle("dim", fontName=font_name, fontSize=8.1, leading=9.2, textColor=colors.HexColor("#27303A"), alignment=TA_LEFT),
        "body": ParagraphStyle("body", fontName=font_name, fontSize=7.8, leading=8.8, textColor=colors.HexColor("#27303A"), alignment=TA_LEFT),
        "risk": ParagraphStyle("risk", fontName=font_name, fontSize=7.8, leading=8.8, textColor=colors.HexColor("#E25C4C"), alignment=TA_LEFT),
    }
    data = [[Paragraph("<b>DIMENSION</b>", styles["header"]), Paragraph("<b>CURRENT ELK + PROMETHEUS</b>", styles["header"]), Paragraph("<b>FULL DYNATRACE ADOPTION</b>", styles["header"])]]
    for dimension, current, dynatrace in ROWS:
        dyn_style = styles["risk"] if dimension in {"Vendor Lock-in", "Migration Risk"} else styles["body"]
        data.append([Paragraph(f"<b>{dimension}</b>", styles["dim"]), Paragraph(current, styles["body"]), Paragraph(f"<b>{dynatrace}</b>" if dimension == "Migration Risk" else dynatrace, dyn_style)])

    table = Table(data, colWidths=[2.15 * inch, 5.05 * inch, 5.15 * inch], rowHeights=[0.45 * inch] + [0.456 * inch] * len(ROWS))
    commands = [
        ("BACKGROUND", (0, 0), (0, 0), colors.HexColor("#162233")),
        ("BACKGROUND", (1, 0), (2, 0), colors.HexColor("#00807D")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 5.8),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5.8),
        ("TOPPADDING", (0, 0), (-1, -1), 1.8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 1.4),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D2DAE0")),
    ]
    for row_idx in range(1, len(ROWS) + 1):
        fill = colors.white if row_idx % 2 else colors.HexColor("#F4F7F9")
        commands.append(("BACKGROUND", (0, row_idx), (-1, row_idx), fill))
    commands.append(("BACKGROUND", (2, len(ROWS)), (2, len(ROWS)), colors.HexColor("#FDECE9")))
    table.setStyle(TableStyle(commands))
    table.wrapOn(pdf, width, height)
    table.drawOn(pdf, 0.48 * inch, height - 1.08 * inch - 5.92 * inch)

    pdf.setFillColor(colors.HexColor("#636F7B"))
    pdf.setFont(font_name, 8)
    pdf.drawString(0.49 * inch, 0.18 * inch, "Assessment focus: capability fit, operating burden, data economics, and migration risk.")
    pdf.save()


def render_preview():
    document = fitz.open(PDF_PATH)
    page = document[0]
    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False)
    pixmap.save(PREVIEW_PATH)
    document.close()


if __name__ == "__main__":
    build_pptx()
    build_pdf()
    render_preview()
    print(PPTX_PATH)
    print(PDF_PATH)
    print(PREVIEW_PATH)
